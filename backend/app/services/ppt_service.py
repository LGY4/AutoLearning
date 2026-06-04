from __future__ import annotations
"""PPT Generation Service — generates learning presentation slides.

Uses python-pptx to create structured presentations from knowledge content.
Supports: title slides, content slides, quiz slides, summary slides.
"""

import io
import logging
from typing import List, Optional
from uuid import UUID

logger = logging.getLogger(__name__)

# ── Color scheme ────────────────────────────────────────────────────
COLORS = {
    "primary": "1A237E",      # Deep blue
    "secondary": "1565C0",    # Blue
    "accent": "2E7D32",       # Green
    "text": "212121",         # Near-black
    "light_text": "757575",   # Gray
    "background": "FFFFFF",   # White
    "highlight": "FF6F00",    # Orange
}


def generate_ppt(
    title: str,
    knowledge_point: str,
    content_sections: List[dict],
    quiz_questions: Optional[List[dict]] = None,
    summary_points: Optional[List[str]] = None,
) -> bytes:
    """Generate a PPT presentation for a learning topic.

    Args:
        title: Presentation title
        knowledge_point: The main knowledge point
        content_sections: List of {heading: str, content: str, bullets: list[str]}
        quiz_questions: Optional quiz questions to include
        summary_points: Optional summary bullet points

    Returns:
        PPT file as bytes
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor.from_string(color or COLORS["text"])
        p.alignment = alignment
        return txBox

    def add_bullet_list(slide, left, top, width, height, items, font_size=16):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor.from_string(COLORS["text"])
            p.space_after = Pt(8)

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    add_text_box(slide, 1, 1.5, 11, 1.5, title, font_size=40, bold=True, color=COLORS["primary"], alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, 1, 3.5, 11, 1, knowledge_point, font_size=24, color=COLORS["secondary"], alignment=PP_ALIGN.CENTER)
    # Decorative line
    from pptx.shapes.connector import Connector
    line = slide.shapes.add_connector(1, Inches(3), Inches(3.2), Inches(10), Inches(3.2))
    line.line.color.rgb = RGBColor.from_string(COLORS["primary"])
    line.line.width = Pt(2)

    # ── Slide 2: Outline ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "目录", font_size=32, bold=True, color=COLORS["primary"])
    outline_items = [s.get("heading", f"第{i+1}部分") for i, s in enumerate(content_sections)]
    if quiz_questions:
        outline_items.append("练习测验")
    if summary_points:
        outline_items.append("总结")
    add_bullet_list(slide, 1, 1.5, 10, 5, outline_items, font_size=20)

    # ── Content slides ──
    for i, section in enumerate(content_sections):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Section number + heading
        add_text_box(slide, 0.5, 0.3, 12, 0.8, f"{i+1}. {section.get('heading', '')}", font_size=28, bold=True, color=COLORS["primary"])

        # Content text
        content = section.get("content", "")
        if content:
            add_text_box(slide, 0.8, 1.5, 11.5, 2.5, content, font_size=16, color=COLORS["text"])

        # Bullet points
        bullets = section.get("bullets", [])
        if bullets:
            add_bullet_list(slide, 0.8, 3.5, 11.5, 3.5, bullets, font_size=16)

        # Page number
        add_text_box(slide, 12, 7, 1, 0.5, str(i + 3), font_size=12, color=COLORS["light_text"], alignment=PP_ALIGN.RIGHT)

    # ── Quiz slides ──
    if quiz_questions:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text_box(slide, 0.5, 0.3, 12, 0.8, "练习测验", font_size=32, bold=True, color=COLORS["primary"])

        for j, q in enumerate(quiz_questions[:5]):
            y_pos = 1.5 + j * 1.2
            if y_pos > 6.5:
                break
            question_text = q.get("question", q.get("stem", ""))
            add_text_box(slide, 0.8, y_pos, 11, 0.5, f"Q{j+1}. {question_text}", font_size=16, bold=True)
            options = q.get("options", [])
            if options:
                opts_text = "    ".join([f"{chr(65+k)}. {v}" for k, v in enumerate(options[:4])])
                add_text_box(slide, 1.2, y_pos + 0.45, 10, 0.5, opts_text, font_size=14, color=COLORS["light_text"])

    # ── Summary slide ──
    if summary_points:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text_box(slide, 0.5, 0.3, 12, 0.8, "总结", font_size=32, bold=True, color=COLORS["primary"])
        add_bullet_list(slide, 1, 1.5, 11, 5, summary_points, font_size=18)

    # ── End slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 1, 2.5, 11, 1.5, "谢谢！", font_size=48, bold=True, color=COLORS["primary"], alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 11, 1, f"学习主题：{knowledge_point}", font_size=20, color=COLORS["secondary"], alignment=PP_ALIGN.CENTER)

    # Export to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def generate_ppt_from_llm(knowledge_point: str, subject: str = "通用", difficulty: str = "medium") -> bytes:
    """Generate a PPT using LLM to create content sections.

    Args:
        knowledge_point: The topic to create a presentation about
        subject: Subject area
        difficulty: easy/medium/hard

    Returns:
        PPT file as bytes
    """
    from app.services.model_gateway import generate_json

    prompt = f"""为知识点"{knowledge_point}"（学科：{subject}，难度：{difficulty}）生成一个教学课件的大纲。

返回严格JSON格式：
{{
    "title": "课件标题",
    "sections": [
        {{"heading": "章节标题", "content": "详细内容说明", "bullets": ["要点1", "要点2", "要点3"]}},
        ...
    ],
    "summary": ["总结要点1", "总结要点2", "总结要点3"],
    "quiz": [
        {{"question": "题目", "options": ["A选项", "B选项", "C选项", "D选项"], "answer": "A选项"}},
        ...
    ]
}}

要求：
- 5-8个章节，由浅入深
- 每章节3-5个要点
- 2-3道练习题
- 内容准确、简洁、适合学习"""

    try:
        result = generate_json(prompt, required_keys=["title", "sections"])
        return generate_ppt(
            title=result.get("title", f"{knowledge_point} 学习课件"),
            knowledge_point=knowledge_point,
            content_sections=result.get("sections", []),
            quiz_questions=result.get("quiz"),
            summary_points=result.get("summary"),
        )
    except Exception as exc:
        logger.error("LLM PPT generation failed: %s", exc)
        # Fallback: basic structure
        return generate_ppt(
            title=f"{knowledge_point} 学习课件",
            knowledge_point=knowledge_point,
            content_sections=[
                {"heading": "概述", "content": f"{knowledge_point}是{subject}中的重要概念。", "bullets": []},
                {"heading": "核心概念", "content": "", "bullets": [f"{knowledge_point}的定义", f"{knowledge_point}的特点", f"{knowledge_point}的应用"]},
                {"heading": "总结", "content": f"掌握{knowledge_point}对于理解{subject}至关重要。", "bullets": []},
            ],
        )
