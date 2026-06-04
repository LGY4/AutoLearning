from __future__ import annotations
"""Enhanced file parser — uses Docling for advanced document parsing.

Falls back to PyPDF2/simple parsing when Docling is unavailable.
Supports: PDF, DOCX, PPTX, XLSX, HTML, Markdown, images (OCR), code files.
"""

import logging
import re
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# ── Docling integration ─────────────────────────────────────────────

_docling_converter = None


def _get_docling_converter():
    """Lazy-load Docling converter (heavy import)."""
    global _docling_converter
    if _docling_converter is not None:
        return _docling_converter
    try:
        from docling.document_converter import DocumentConverter
        _docling_converter = DocumentConverter()
        logger.info("Docling document converter initialized")
    except Exception as exc:
        logger.warning("Docling unavailable, falling back to basic parsing: %s", exc)
        _docling_converter = False  # sentinel: tried and failed
    return _docling_converter


def _parse_with_docling(filename: str, content: bytes) -> Optional[str]:
    """Parse document using Docling (supports PDF, DOCX, PPTX, HTML, images)."""
    converter = _get_docling_converter()
    if converter is False or converter is None:
        return None

    try:
        import tempfile
        import os

        ext = Path(filename).suffix.lower()
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            result = converter.convert(tmp_path)
            doc = result.document

            # Export to markdown — preserves structure, tables, headings
            md_content = doc.export_to_markdown()

            # Also extract tables as structured data
            tables_text = ""
            try:
                for table in doc.tables:
                    table_md = table.export_to_markdown()
                    if table_md.strip():
                        tables_text += f"\n\n[表格]\n{table_md}\n"
            except Exception:
                pass

            return md_content + tables_text
        finally:
            os.unlink(tmp_path)
    except Exception as exc:
        logger.warning("Docling parsing failed for %s: %s", filename, exc)
        return None


# ── Basic parsers (fallback) ────────────────────────────────────────

def parse_markdown(content: str) -> str:
    """Extract structured text from Markdown content."""
    lines = content.split("\n")
    result: List[str] = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            result.append(stripped)
        else:
            cleaned = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', stripped)
            cleaned = re.sub(r'[*_`~]+', '', cleaned)
            cleaned = re.sub(r'^[-*+]\s+', '• ', cleaned)
            cleaned = re.sub(r'^\d+\.\s+', '', cleaned)
            if cleaned:
                result.append(cleaned)
    return "\n".join(result)


def parse_pdf_bytes(data: bytes) -> str:
    """Extract text from PDF bytes using PyPDF2 (fallback)."""
    try:
        import io
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages: List[str] = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text.strip())
        return "\n\n".join(pages)
    except ImportError:
        try:
            import io
            import pdfplumber
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                pages = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        pages.append(text.strip())
                return "\n\n".join(pages)
        except ImportError:
            raise RuntimeError("No PDF parser installed. pip install PyPDF2")


def parse_docx_bytes(data: bytes) -> str:
    """Extract text from DOCX bytes."""
    try:
        import io
        from docx import Document
        doc = Document(io.BytesIO(data))
        parts: List[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure
                if para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    parts.append(f"{'#' * int(level)} {text}")
                else:
                    parts.append(text)
        # Extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n[表格]\n" + "\n".join(rows))
        return "\n".join(parts)
    except ImportError:
        raise RuntimeError("python-docx not installed. pip install python-docx")


def parse_pptx_bytes(data: bytes) -> str:
    """Extract text from PPTX bytes."""
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text: List[str] = [f"## 幻灯片 {i}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        slide_text.append("[表格]\n" + "\n".join(rows))
            if len(slide_text) > 1:
                parts.append("\n".join(slide_text))
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("python-pptx not installed. pip install python-pptx")


def parse_xlsx_bytes(data: bytes) -> str:
    """Extract text from XLSX bytes."""
    try:
        import io
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: List[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows[:100]))  # limit rows
        wb.close()
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("openpyxl not installed. pip install openpyxl")


# ── Main entry point ────────────────────────────────────────────────

def parse_uploaded_file(filename: str, content: bytes) -> str:
    """Parse an uploaded file. Uses Docling when available, falls back to basic parsers."""
    ext = Path(filename).suffix.lower()

    # Docling handles: .pdf, .docx, .pptx, .html, .htm, .png, .jpg, .jpeg, .tiff
    docling_extensions = {".pdf", ".docx", ".pptx", ".html", ".htm", ".png", ".jpg", ".jpeg", ".tiff"}
    if ext in docling_extensions:
        result = _parse_with_docling(filename, content)
        if result:
            return result
        # Fall through to basic parsers

    # Basic parsers
    if ext in (".md", ".markdown"):
        return parse_markdown(content.decode("utf-8", errors="replace"))

    elif ext == ".pdf":
        return parse_pdf_bytes(content)

    elif ext in (".docx",):
        return parse_docx_bytes(content)

    elif ext in (".pptx",):
        return parse_pptx_bytes(content)

    elif ext in (".xlsx", ".xls"):
        return parse_xlsx_bytes(content)

    elif ext in (".txt", ".text", ".csv"):
        return content.decode("utf-8", errors="replace")

    elif ext in (".json",):
        import json
        data = json.loads(content.decode("utf-8"))
        if isinstance(data, list):
            parts = []
            for item in data:
                title = item.get("title", "")
                body = item.get("content", item.get("description", ""))
                if title or body:
                    parts.append(f"## {title}\n{body}" if title else body)
            return "\n\n".join(parts)
        return json.dumps(data, ensure_ascii=False, indent=2)

    else:
        # Try UTF-8 text as last resort
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            raise ValueError(f"Unsupported file type: {ext}")


def get_supported_extensions() -> List[str]:
    """Return list of supported file extensions."""
    return [
        ".pdf", ".docx", ".pptx", ".xlsx", ".html", ".htm",
        ".md", ".markdown", ".txt", ".text", ".csv", ".json",
        ".png", ".jpg", ".jpeg", ".tiff",
        ".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs",
    ]
